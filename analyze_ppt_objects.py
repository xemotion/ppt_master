from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.chart import Chart
from pptx.table import Table
from pptx.enum.shapes import MSO_SHAPE
import json
import os

def emu_to_inches(emu):
    """EMU를 인치로 변환합니다."""
    return round(emu / 914400, 2)

def emu_to_pixels(emu, dpi=96):
    """EMU를 픽셀로 변환합니다."""
    return round((emu / 914400) * dpi, 2)

def get_shape_type_name(shape):
    """shape의 타입 이름을 반환합니다."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return "텍스트 상자"
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "이미지"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "차트"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "표"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if hasattr(shape, "auto_shape_type"):
                return f"도형 ({MSO_SHAPE(shape.auto_shape_type).name})"
            return "도형"
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "플레이스홀더"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP_SHAPE:
            return "그룹화된 도형"
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return "선"
        else:
            return "기타"
    except:
        return "알 수 없음"

def get_shape_description(shape, text_content):
    """텍스트 내용을 기반으로 객체의 역할을 설명합니다."""
    # 위치 기반 설명
    position_desc = ""
    if shape.top < 1000000:  # 상단에 가까운 경우
        position_desc = "상단에 위치한 "
    elif shape.top > 5000000:  # 하단에 가까운 경우
        position_desc = "하단에 위치한 "
    
    # 크기 기반 설명
    size_desc = ""
    if shape.width > 5000000:  # 큰 크기
        size_desc = "큰 "
    elif shape.width < 2000000:  # 작은 크기
        size_desc = "작은 "
    
    # 텍스트 길이 기반 설명
    text_length = len(text_content)
    if text_length > 50:
        text_desc = "긴 텍스트를 포함한 "
    elif text_length < 10:
        text_desc = "짧은 텍스트를 포함한 "
    else:
        text_desc = ""
    
    # 텍스트 내용 기반 설명
    content_desc = ""
    if "제목" in text_content or "Title" in text_content:
        content_desc = "제목"
    elif "부제목" in text_content or "Subtitle" in text_content:
        content_desc = "부제목"
    elif "목차" in text_content or "Contents" in text_content:
        content_desc = "목차"
    elif "참고" in text_content or "Reference" in text_content:
        content_desc = "참고 문헌"
    elif "결론" in text_content or "Conclusion" in text_content:
        content_desc = "결론"
    elif "소개" in text_content or "Introduction" in text_content:
        content_desc = "소개"
    else:
        content_desc = "본문"
    
    # 최종 설명 조합
    description = f"{position_desc}{size_desc}{text_desc}{content_desc}을 위한 객체"
    
    return description

def get_shape_identifier(shape, slide_width, slide_height):
    """객체의 식별 정보를 추출합니다."""
    # EMU를 인치로 변환
    left_inches = emu_to_inches(shape.left)
    top_inches = emu_to_inches(shape.top)
    width_inches = emu_to_inches(shape.width)
    height_inches = emu_to_inches(shape.height)
    
    # EMU를 픽셀로 변환
    left_pixels = emu_to_pixels(shape.left)
    top_pixels = emu_to_pixels(shape.top)
    width_pixels = emu_to_pixels(shape.width)
    height_pixels = emu_to_pixels(shape.height)
    
    identifier = {
        "position": {
            "emu": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            },
            "inches": {
                "left": left_inches,
                "top": top_inches,
                "width": width_inches,
                "height": height_inches
            },
            "pixels": {
                "left": left_pixels,
                "top": top_pixels,
                "width": width_pixels,
                "height": height_pixels
            },
            "relative_position": {
                "left_percent": round((shape.left / slide_width) * 100, 2),
                "top_percent": round((shape.top / slide_height) * 100, 2),
                "width_percent": round((shape.width / slide_width) * 100, 2),
                "height_percent": round((shape.height / slide_height) * 100, 2)
            }
        }
    }
    
    # 텍스트 기반 식별자
    try:
        if hasattr(shape, "text") and shape.text.strip():
            text_content = shape.text.strip()
            identifier["text_content"] = text_content
            identifier["text_length"] = len(text_content)
            identifier["description"] = get_shape_description(shape, text_content)
            
            # 텍스트 스타일 정보
            if hasattr(shape, "text_frame"):
                if shape.text_frame.paragraphs:
                    first_para = shape.text_frame.paragraphs[0]
                    if first_para.runs:
                        first_run = first_para.runs[0]
                        identifier["text_style"] = {
                            "font_size": first_run.font.size.pt if first_run.font.size else None,
                            "font_name": first_run.font.name if first_run.font.name else None,
                            "is_bold": first_run.font.bold,
                            "is_italic": first_run.font.italic
                        }
    except:
        pass
        
    # 도형 타입 기반 식별자
    try:
        if hasattr(shape, "auto_shape_type") and shape.auto_shape_type is not None:
            identifier["shape_type"] = MSO_SHAPE(shape.auto_shape_type).name
    except:
        pass
        
    # 이미지 식별자
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            identifier["image_type"] = "picture"
            if hasattr(shape, "image"):
                identifier["image_size"] = {
                    "width": shape.image.width,
                    "height": shape.image.height
                }
    except:
        pass
        
    # 차트 식별자
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart = shape.chart
            identifier["chart_type"] = str(chart.chart_type)
            if chart.chart_title:
                identifier["chart_title"] = chart.chart_title.text
    except:
        pass
        
    # 표 식별자
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            identifier["table_dimensions"] = {
                "rows": table.rows.__len__(),
                "columns": table.columns.__len__()
            }
            # 표의 텍스트 내용도 추가
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            identifier["table_content"] = table_data
    except:
        pass
        
    # 스타일 정보
    try:
        if hasattr(shape, "fill"):
            if shape.fill.type is not None:
                identifier["fill_type"] = str(shape.fill.type)
    except:
        pass
        
    try:
        if hasattr(shape, "line"):
            if shape.line.color.rgb is not None:
                identifier["line_color"] = f"#{shape.line.color.rgb:06x}"
    except:
        pass
    
    return identifier

def analyze_shape(shape, slide_width, slide_height):
    """각 shape의 세부 정보를 분석합니다."""
    try:
        shape_info = {
            "name": shape.name,
            "type": shape.shape_type,
            "type_name": MSO_SHAPE_TYPE(shape.shape_type).name if shape.shape_type in MSO_SHAPE_TYPE._value2member_map_ else "UNKNOWN",
            "type_description": get_shape_type_name(shape),
            "id": shape.shape_id,
            "identifier": get_shape_identifier(shape, slide_width, slide_height)
        }
        return shape_info
    except Exception as e:
        return {
            "error": str(e),
            "type": "unknown"
        }

def analyze_pptx(pptx_path):
    """PPT 파일의 모든 슬라이드와 객체들을 분석합니다."""
    prs = Presentation(pptx_path)
    result = {
        "total_slides": len(prs.slides),
        "slides": []
    }
    
    for slide_number, slide in enumerate(prs.slides, 1):
        # 슬라이드 크기 정보
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        slide_info = {
            "slide_number": slide_number,
            "shapes": [],
            "layout_name": slide.slide_layout.name,
            "total_shapes": len(slide.shapes),
            "slide_dimensions": {
                "emu": {
                    "width": slide_width,
                    "height": slide_height
                },
                "inches": {
                    "width": emu_to_inches(slide_width),
                    "height": emu_to_inches(slide_height)
                },
                "pixels": {
                    "width": emu_to_pixels(slide_width),
                    "height": emu_to_pixels(slide_height)
                }
            }
        }
        
        for shape in slide.shapes:
            shape_info = analyze_shape(shape, slide_width, slide_height)
            # 텍스트가 있는 객체만 추가
            if "text_content" in shape_info.get("identifier", {}):
                slide_info["shapes"].append(shape_info)
            
        result["slides"].append(slide_info)
    
    return result

def save_to_json(data, output_path):
    """분석 결과를 JSON 파일로 저장합니다."""
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def main():
    pptx_path = "test.pptx"
    output_path = "ppt_analysis.json"  # ppt_master 디렉토리에 저장
    
    try:
        result = analyze_pptx(pptx_path)
        save_to_json(result, output_path)
        print(f"분석 결과가 {output_path}에 저장되었습니다.")
    except Exception as e:
        print(f"에러 발생: {str(e)}")

if __name__ == "__main__":
    main() 