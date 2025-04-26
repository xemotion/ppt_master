from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import os

def generate_object_schema(pptx_path):
    """PPT 파일을 분석하여 객체 스키마를 생성합니다."""
    prs = Presentation(pptx_path)
    schema = {
        "total_slides": len(prs.slides),
        "slides": []
    }
    
    for slide_number, slide in enumerate(prs.slides, 1):
        slide_schema = {
            "slide_number": slide_number,
            "layout_name": slide.slide_layout.name,
            "objects": []
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                object_schema = {
                    "type": "text_object",
                    "content": shape.text.strip(),
                    "position": {
                        "top": shape.top,
                        "left": shape.left,
                        "width": shape.width,
                        "height": shape.height
                    },
                    "description": generate_description(shape)
                }
                slide_schema["objects"].append(object_schema)
        
        schema["slides"].append(slide_schema)
    
    return schema

def generate_description(shape):
    """객체의 설명을 생성합니다."""
    # 위치 기반 설명
    position_desc = ""
    if shape.top < 1000000:  # 상단에 가까운 경우
        position_desc = "상단에 위치한 "
    elif shape.top > 5000000:  # 하단에 가까운 경우
        position_desc = "하단에 위치한 "
    
    # 크기 기반 설명
    size_desc = ""
    if shape.width > 5000000:  # 큰 크기
        size_desc = "큰 "
    elif shape.width < 2000000:  # 작은 크기
        size_desc = "작은 "
    
    # 텍스트 내용 기반 설명
    text_content = shape.text.strip()
    content_desc = ""
    if "제목" in text_content or "Title" in text_content:
        content_desc = "제목"
    elif "부제목" in text_content or "Subtitle" in text_content:
        content_desc = "부제목"
    elif "목차" in text_content or "Contents" in text_content:
        content_desc = "목차"
    elif "참고" in text_content or "Reference" in text_content:
        content_desc = "참고 문헌"
    elif "결론" in text_content or "Conclusion" in text_content:
        content_desc = "결론"
    elif "소개" in text_content or "Introduction" in text_content:
        content_desc = "소개"
    else:
        content_desc = "본문"
    
    # 최종 설명 조합
    description = f"{position_desc}{size_desc}{content_desc}을 위한 객체"
    
    return description

def save_schema(schema, output_path):
    """스키마를 JSON 파일로 저장합니다."""
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(schema, f, ensure_ascii=False, indent=2)

def main():
    pptx_path = "test.pptx"
    output_path = "ppt_object_schema.json"
    
    try:
        schema = generate_object_schema(pptx_path)
        save_schema(schema, output_path)
        print(f"객체 스키마가 {output_path}에 저장되었습니다.")
    except Exception as e:
        print(f"에러 발생: {str(e)}")

if __name__ == "__main__":
    main() 