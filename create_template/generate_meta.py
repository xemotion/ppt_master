import json
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dotenv import load_dotenv
from openai import AzureOpenAI
from typing import Dict, Any, List, Optional
from pathlib import Path
import logging

# 공통 모듈에서 함수 가져오기
from ppt_common import (
    get_shape_position,
    get_type_info,
    make_element_id,
    logger,
    generate_position_key,
    is_tag_identifier,
    is_special_content
)

# 로깅 설정
logging.basicConfig(
    filename="ppt_process.log",
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def determine_element_role(shape, pos, slide_number, slide_width, slide_height, shape_type):
    """
    슬라이드 내 요소의 위치와 유형을 고려하여 의미있는 역할 이름을 생성합니다.
    
    Args:
        shape: PPT 도형 객체
        pos: 위치 정보 딕셔너리 (left_percent, top_percent, width_percent, height_percent)
        slide_number: 슬라이드 번호 (함수 내부에서는 사용하지 않음)
        slide_width: 슬라이드 너비
        slide_height: 슬라이드 높이
        shape_type: 도형 유형 (TEXT_BOX, AUTO_SHAPE 등)
        
    Returns:
        str: 요소의 역할 이름
    """
    # 위치 결정 (상단, 중간, 하단)
    if pos["top_percent"] < 20:
        vertical_pos = "top"
    elif pos["top_percent"] > 70:
        vertical_pos = "bottom"
    else:
        vertical_pos = "middle"
    
    # 위치 결정 (좌측, 중앙, 우측)
    if pos["left_percent"] < 30:
        horizontal_pos = "left"
    elif pos["left_percent"] > 60:
        horizontal_pos = "right"
    else:
        horizontal_pos = "center"
    
    # 크기 결정 (크기에 따라 중요도 부여)
    if pos["width_percent"] > 70 or pos["height_percent"] > 30:
        size = "main"
    elif pos["width_percent"] > 40 or pos["height_percent"] > 15:
        size = "sub"
    else:
        size = "detail"
    
    # 요소 유형에 따른 기능적 역할
    if shape_type == "TEXT_BOX":
        # 텍스트 크기와 위치로 역할 추론
        if vertical_pos == "top" and (size == "main" or size == "sub"):
            functional_role = "title"
        elif vertical_pos == "top" and size == "detail":
            functional_role = "header"
        elif vertical_pos == "bottom":
            functional_role = "footer"
        else:
            functional_role = "content"
    elif shape_type == "AUTO_SHAPE":
        functional_role = "shape"
    elif shape_type == "PICTURE":
        functional_role = "image"
    elif shape_type == "GROUP":
        functional_role = "group"
    elif shape_type == "TABLE":
        functional_role = "table"
    elif shape_type == "CHART":
        functional_role = "chart"
    else:
        functional_role = "element"
    
    # 특수 케이스 처리: 상단 제목
    if vertical_pos == "top" and horizontal_pos == "center" and size == "main" and shape_type == "TEXT_BOX":
        return f"{vertical_pos}_{horizontal_pos}_main_title"
    
    # 일반적인 이름 생성 (슬라이드 번호 제외)
    return f"{vertical_pos}_{horizontal_pos}_{size}_{functional_role}"

def get_position_category(pos):
    """위치에 따른 카테고리를 더 세분화하여 반환합니다."""
    # 수직 위치를 5개 구간으로 나눔
    if pos["top_percent"] < 20:
        vertical = "top"
    elif pos["top_percent"] < 40:
        vertical = "upper"
    elif pos["top_percent"] < 60:
        vertical = "middle"
    elif pos["top_percent"] < 80:
        vertical = "lower"
    else:
        vertical = "bottom"
    
    # 수평 위치를 5개 구간으로 나눔
    if pos["left_percent"] < 20:
        horizontal = "far_left"
    elif pos["left_percent"] < 40:
        horizontal = "left"
    elif pos["left_percent"] < 60:
        horizontal = "center"
    elif pos["left_percent"] < 80:
        horizontal = "right"
    else:
        horizontal = "far_right"
    
    return vertical, horizontal

def call_llm_for_meta(text, pos, type_name, slide_number, client, deployment_name):
    """
    텍스트 요소에 대한 역할과 설명을 생성합니다.
    LLM API를 호출하거나 규칙 기반으로 역할을 생성합니다.
    """
    try:
        # 위치 카테고리 가져오기
        vertical, horizontal = get_position_category(pos)
        
        # 위치별 카운터 초기화 (없으면 생성)
        if not hasattr(call_llm_for_meta, 'position_counters'):
            call_llm_for_meta.position_counters = {}
        
        position_key = f"{vertical}_{horizontal}"
        if position_key not in call_llm_for_meta.position_counters:
            call_llm_for_meta.position_counters[position_key] = {}
        
        # 태그/라벨 확인
        if is_tag_identifier(text):
            print(f"태그/라벨 요소 감지: '{text}'")
            
            # 태그 카운터 증가
            if 'tag' not in call_llm_for_meta.position_counters[position_key]:
                call_llm_for_meta.position_counters[position_key]['tag'] = 0
            call_llm_for_meta.position_counters[position_key]['tag'] += 1
            
            tag_count = call_llm_for_meta.position_counters[position_key]['tag']
            role = f"{vertical}_{horizontal}_tag_{tag_count}"
            description = f"슬라이드 {slide_number}의 태그 요소 (텍스트: '{text}')"
            
            return role, description
                
        # 숫자나 특수 기호로만 된 내용 체크
        if is_special_content(text):
            print(f"특수 콘텐츠 감지: '{text}'")
            
            # 특수 콘텐츠는 role을 원본 텍스트로 사용
            role = text
            description = f"슬라이드 {slide_number}의 특수 콘텐츠 (원본: '{text}', 위치: {pos['left_percent']:.1f}%, {pos['top_percent']:.1f}%)"
            return role, description
        
        # LLM API 호출 여부 확인
        use_llm = True  # LLM 호출 활성화
        
        if use_llm and all([client, deployment_name]):
            system_prompt = f"""
                You are an expert in slide design and template structure analysis.

                You will receive a text element from a presentation slide along with its position and shape information.

                Your task is to:
                1. Assign a descriptive role name (`role`) based on its function and type. 
                Important rules for role names:
                - Format: [vertical]_[horizontal]_[type]_[index]
                - Vertical position: {vertical}
                - Horizontal position: {horizontal}
                - Type: 'title', 'content', 'note', 'list', 'special'
                - Index: Will be added automatically
                - DO NOT include any position numbers or coordinates
                - GOOD examples: 'upper_left_title', 'middle_center_content'
                
                2. Write a structural `description` that clearly explains the **role and functional purpose of this element in the overall slide template**.
                - DO NOT mention the actual content of the text.
                - DO NOT describe the topic (e.g. SW, DB, AI).
                - Describe only the visual structure, layout function, importance level, and placement logic.
                - Example: "This element is placed at the center top of the slide and serves as the main heading, establishing the visual hierarchy of the architecture overview template."

                Current element information:
                - Text: {text}
                - Position: {pos}
                - Type: {type_name}
                - Slide Number: {slide_number}

                Respond with a JSON object:
                {{
                "role": "...",
                "description": "..."
                }}
                Respond ONLY with a valid JSON. Do not include any markdown, bullet points, or extra commentary.
                """.strip()

            user_prompt = {
                "text": text,
                "position": pos,
                "type": type_name,
                "slide_number": slide_number
            }
            
            try:
                print(f"LLM API 호출 중... (슬라이드 {slide_number}, 텍스트: '{text[:30]}{'...' if len(text) > 30 else ''}')")
                response = client.chat.completions.create(
                    model=deployment_name,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": json.dumps(user_prompt, ensure_ascii=False, indent=2)}
                    ],
                    temperature=0.3,  # 더 일관된 결과를 위해 temperature 낮춤
                    max_tokens=300
                )
                response_text = response.choices[0].message.content.strip()
                print(f"LLM Raw Response (슬라이드 {slide_number}):", response_text)

                try:
                    # 마크다운 코드 블록 제거
                    clean_response = response_text.strip()
                    if clean_response.startswith('```'):
                        # 첫 번째 ``` 이후의 텍스트 추출
                        clean_response = clean_response.split('\n', 1)[1]
                    if clean_response.endswith('```'):
                        # 마지막 ``` 제거
                        clean_response = clean_response.rsplit('\n', 1)[0]
                    # json이나 다른 언어 지정자 제거
                    if clean_response.startswith('{'):
                        data = json.loads(clean_response)
                    else:
                        # 첫 번째 { 찾기
                        json_start = clean_response.find('{')
                        if json_start != -1:
                            clean_response = clean_response[json_start:]
                            data = json.loads(clean_response)
                        else:
                            raise ValueError("No JSON object found in response")
                            
                    base_role = data["role"]
                    description = data["description"]
                    
                    # role 타입 추출 및 카운터 증가
                    role_type = base_role.split('_')[-1]  # 마지막 부분을 타입으로 사용
                    if role_type not in call_llm_for_meta.position_counters[position_key]:
                        call_llm_for_meta.position_counters[position_key][role_type] = 0
                    call_llm_for_meta.position_counters[position_key][role_type] += 1
                    
                    # 최종 role 생성 (카운터 포함)
                    role = f"{base_role}_{call_llm_for_meta.position_counters[position_key][role_type]}"
                    
                    return role, description
                except Exception as e:
                    print(f"[Parse Error] {e}")
                    print(f"[RAW Response] {response_text}")
                    
                    # 특수 문자나 숫자만 있는 경우 원본 텍스트를 role로 사용
                    if is_special_content(text):
                        return text, f"슬라이드 {slide_number}의 특수 콘텐츠 (원본: '{text}', 위치: {pos['left_percent']:.1f}%, {pos['top_percent']:.1f}%)"
                    
                    # 일반 텍스트의 경우 기본 role 생성
                    if 'content' not in call_llm_for_meta.position_counters[position_key]:
                        call_llm_for_meta.position_counters[position_key]['content'] = 0
                    call_llm_for_meta.position_counters[position_key]['content'] += 1
                    
                    default_role = f"{vertical}_{horizontal}_content_{call_llm_for_meta.position_counters[position_key]['content']}"
                    return default_role, f"JSON 파싱 오류로 인한 기본 역할 생성: {str(e)[:100]}"
            except Exception as e:
                print(f"[LLM API Error] {e}")
                # API 호출 실패 시 기본 role 생성
                if 'content' not in call_llm_for_meta.position_counters[position_key]:
                    call_llm_for_meta.position_counters[position_key]['content'] = 0
                call_llm_for_meta.position_counters[position_key]['content'] += 1
                
                default_role = f"{vertical}_{horizontal}_content_{call_llm_for_meta.position_counters[position_key]['content']}"
                return default_role, f"LLM API 오류로 인한 기본 역할 생성: {str(e)[:100]}"
        
        # LLM을 사용하지 않거나 호출에 실패한 경우 의미 있는 role 생성
        size_desc = "main" if pos["width_percent"] > 70 or pos["height_percent"] > 30 else "sub" if pos["width_percent"] > 40 or pos["height_percent"] > 15 else "detail"
        
        type_desc = {
            "TEXT_BOX": "content",
            "AUTO_SHAPE": "shape",
            "PICTURE": "image",
            "GROUP": "group",
            "TABLE": "table",
            "CHART": "chart"
        }.get(type_name, "element")
        
        # 타입별 카운터 증가
        if type_desc not in call_llm_for_meta.position_counters[position_key]:
            call_llm_for_meta.position_counters[position_key][type_desc] = 0
        call_llm_for_meta.position_counters[position_key][type_desc] += 1
        
        # 최종 role 생성
        role_name = f"{vertical}_{horizontal}_{size_desc}_{type_desc}_{call_llm_for_meta.position_counters[position_key][type_desc]}"
        description = f"슬라이드 {slide_number}의 {vertical} {horizontal}에 위치한 {size_desc} {type_desc} 요소"
        
        return role_name, description
    except Exception as e:
        # 예상치 못한 오류 발생 시 안전하게 처리
        print(f"[Unexpected Error in call_llm_for_meta] {e}")
        if 'error' not in call_llm_for_meta.position_counters[position_key]:
            call_llm_for_meta.position_counters[position_key]['error'] = 0
        call_llm_for_meta.position_counters[position_key]['error'] += 1
        
        default_role = f"slide_{slide_number}_element_{call_llm_for_meta.position_counters[position_key]['error']}"
        return default_role, f"오류: {str(e)[:100]}"

# 위치 카운터 초기화
call_llm_for_meta.position_counters = {}

def generate_unique_key(text: str, position: dict) -> str:
    """텍스트와 위치 정보를 조합하여 고유 키 생성"""
    # 위치를 5% 단위로 반올림하여 근접 위치는 같은 키로 처리
    left_group = round(position['left_percent'] / 5) * 5
    top_group = round(position['top_percent'] / 5) * 5
    return f"{text}__pos_{left_group}_{top_group}"

def extract_text_from_shape(shape) -> dict:
    """도형에서 텍스트와 폰트 컬러를 추출합니다."""
    if not hasattr(shape, "text_frame"):
        return {"text": "", "font_color": (0, 0, 0)}
    
    text_parts = []
    font_colors = []
    
    for paragraph in shape.text_frame.paragraphs:
        paragraph_runs = []
        paragraph_color = None
        
        for run in paragraph.runs:
            if run.text and run.text.strip():
                paragraph_runs.append(run.text.strip())
                # 폰트 컬러 추출
                try:
                    if hasattr(run.font, 'color') and run.font.color and hasattr(run.font.color, 'rgb'):
                        paragraph_color = run.font.color.rgb
                except AttributeError:
                    pass
        
        if paragraph_runs:
            text_parts.append(" ".join(paragraph_runs))
            font_colors.append(paragraph_color or (0, 0, 0))  # 기본값은 검정색
    
    # 모든 paragraph의 텍스트를 합치고, 첫 번째 유효한 폰트 컬러 사용
    return {
        "text": "\n".join(text_parts),
        "font_color": font_colors[0] if font_colors else (0, 0, 0)
    }

def extract_table_info(shape) -> Optional[Dict[str, Any]]:
    """표에서 정보를 추출합니다."""
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        return None
        
    table_data = []
    for row_idx, row in enumerate(shape.table.rows, 1):
        row_data = []
        for col_idx, cell in enumerate(row.cells, 1):
            cell_text = "\n".join([p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip()])
            if cell_text:
                row_data.append({
                    "text": cell_text,
                    "row": row_idx,
                    "col": col_idx
                })
        if row_data:
            table_data.extend(row_data)
            
    return table_data if table_data else None

def process_shape(shape, slide_number: int, slide_width: int, slide_height: int) -> Optional[Dict[str, Any]]:
    """개별 도형을 처리하고 메타 정보를 추출합니다."""
    # 기본 위치 정보 추출
    pos = get_shape_position(shape, slide_width, slide_height)
    type_name = get_type_info(shape)
    
    # 표 처리
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table_info = extract_table_info(shape)
        if table_info:
            return {
                "type": "table",
                "position": pos,
                "cells": table_info
            }
        return None
    
    # 텍스트 추출
    text_info = extract_text_from_shape(shape)
    if not text_info["text"]:
        return None
        
    # 위치 기반 키 생성
    position_key = generate_position_key(text_info["text"], pos)
    
    # 메타 정보 구성
    return {
        "id": make_element_id(slide_number, pos, type_name),
        "type": type_name,
        "text": text_info["text"],
        "position": pos,
        "position_key": position_key,
        "font_color": text_info["font_color"]
    }

def extract_meta_info(slide, slide_number: int, slide_width: int, slide_height: int) -> Dict[str, Any]:
    """단일 슬라이드에서 메타 정보를 추출합니다."""
    meta_info = {
        "slide_number": slide_number,
        "slide_width": slide_width,
        "slide_height": slide_height,
        "fields": {}
    }

    # OpenAI 클라이언트 설정
    load_dotenv()
    client = AzureOpenAI(
        api_key=os.getenv("AZURE_OPENAI_API_KEY"),
        api_version=os.getenv("AZURE_OPENAI_VERSION"),
        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT")
    )
    deployment_name = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME")

    # 모든 도형 처리
    text_counts = {}  # 텍스트 중복 카운트용
    
    for shape in slide.shapes:
        # 도형의 텍스트와 폰트 컬러 추출
        if not hasattr(shape, "text_frame"):
            continue
            
        text_info = extract_text_from_shape(shape)
        if not text_info["text"]:
            continue
            
        # 도형의 위치 정보 추출
        pos = get_shape_position(shape, slide_width, slide_height)
        position_key = generate_position_key(text_info["text"], pos)
        
        # 텍스트 중복 처리
        if text_info["text"] in text_counts:
            text_counts[text_info["text"]] += 1
            position_key = f"{position_key}_{text_counts[text_info['text']]}"
        else:
            text_counts[text_info["text"]] = 1
        
        # 도형 유형 확인
        type_name = MSO_SHAPE_TYPE(shape.shape_type).name if shape.shape_type in MSO_SHAPE_TYPE._value2member_map_ else "UNKNOWN"
        
        # role과 description 생성
        role, description = call_llm_for_meta(
            text=text_info["text"],
            pos=pos,
            type_name=type_name,
            slide_number=slide_number,
            client=client,
            deployment_name=deployment_name
        )
        
        # 위치 기반 키로 저장
        meta_info["fields"][position_key] = {
            "type": type_name,
            "original_text": text_info["text"],
            "position": pos,
            "element_id": f"element_slide{slide_number}_l{int(pos['left_percent'])}_t{int(pos['top_percent'])}",
            "role": role,
            "role_description": description,
            "font_color": text_info["font_color"]  # 폰트 컬러 저장
        }
        
        # 태그/라벨 요소 처리
        if is_tag_identifier(text_info["text"]):
            meta_info["fields"][position_key]["is_tag"] = True

    return meta_info

def process_meta_info(prs, slide_number, slide, slide_width, slide_height):
    """
    단일 슬라이드의 메타 정보를 추출하고 처리합니다.
    """
    try:
        # 메타 정보 추출
        meta_info = extract_meta_info(slide, slide_number, slide_width, slide_height)
        return meta_info
    except Exception as e:
        logger.error(f"슬라이드 {slide_number} 메타 정보 추출 중 오류: {str(e)}")
        raise

def save_meta_info(meta_info, output_path):
    """
    메타 정보를 JSON 파일로 저장합니다.
    """
    try:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(meta_info, f, ensure_ascii=False, indent=2)
        logger.info(f"메타 정보 저장 완료: {output_path}")
    except Exception as e:
        logger.error(f"메타 정보 저장 중 오류: {str(e)}")
        raise 
