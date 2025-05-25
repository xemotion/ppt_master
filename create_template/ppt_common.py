import re
import logging
from typing import Dict, List, Optional, Union, Any, Tuple
import os
import json
from pathlib import Path
from pptx import Presentation

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 로깅 설정
logging.basicConfig(
    filename="ppt_process.log", 
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


def safe_get_font_color(run):
    """
    안전하게 폰트 컬러를 추출합니다.
    
    Args:
        run: 텍스트 run 객체
        
    Returns:
        tuple or None: (r,g,b) 컬러값 또는 None
    """
    if not run.text or not run.text.strip():
        return None

    try:
        font = run.font
        if not font:
            return None
            
        color = getattr(font, "color", None)
        if not color:
            return None

        # RGB 값이 직접 있는 경우
        rgb = getattr(color, "rgb", None)
        if rgb:
            try:
                # RGB 값을 튜플로 변환
                if hasattr(rgb, '_value'):  # RGBColor 객체인 경우
                    hex_str = f"{rgb._value:06x}"
                    r = int(hex_str[:2], 16)
                    g = int(hex_str[2:4], 16)
                    b = int(hex_str[4:], 16)
                    return (r, g, b)
                elif isinstance(rgb, (tuple, list)) and len(rgb) == 3:  # RGB 튜플인 경우
                    return tuple(rgb)
                elif hasattr(rgb, 'r') and hasattr(rgb, 'g') and hasattr(rgb, 'b'):  # RGB 속성이 있는 객체
                    return (rgb.r, rgb.g, rgb.b)
            except Exception as e:
                logger.warning(f"RGB 값 변환 중 오류: {e}")
                return None

        # theme color나 다른 색상 형식이 있는 경우 처리
        theme_color = getattr(color, "theme_color", None)
        if theme_color:
            # theme color를 RGB로 변환하는 로직 추가 가능
            pass

    except Exception as e:
        logger.warning(f"폰트 컬러 추출 중 오류: {e}")
        return None

    return None


#################################################
# 공통 유틸리티 함수
#################################################

def get_shape_position(shape, slide_width, slide_height):
    """
    슬라이드 내의 도형 위치를 백분율로 반환합니다.
    
    Args:
        shape: PowerPoint 도형 객체
        slide_width: 슬라이드 너비 (EMU)
        slide_height: 슬라이드 높이 (EMU)
        
    Returns:
        dict: 도형의 위치 정보 (백분율)
    """
    # EMU 값을 정수로 변환 (._value 속성 사용)
    width = slide_width._value if hasattr(slide_width, '_value') else int(slide_width)
    height = slide_height._value if hasattr(slide_height, '_value') else int(slide_height)
    left = shape.left._value if hasattr(shape.left, '_value') else int(shape.left)
    top = shape.top._value if hasattr(shape.top, '_value') else int(shape.top)
    shape_width = shape.width._value if hasattr(shape.width, '_value') else int(shape.width)
    shape_height = shape.height._value if hasattr(shape.height, '_value') else int(shape.height)
    
    return {
        "left_percent": round((left / width) * 100, 2),
        "top_percent": round((top / height) * 100, 2),
        "width_percent": round((shape_width / width) * 100, 2),
        "height_percent": round((shape_height / height) * 100, 2)
    }

def get_type_info(shape):
    """도형의 유형 이름을 반환합니다."""
    return MSO_SHAPE_TYPE(shape.shape_type).name

def make_element_id(slide_number, pos, type_name):
    """요소의 고유 ID를 생성합니다."""
    return f"element_slide{slide_number}_l{int(pos['left_percent'])}_t{int(pos['top_percent'])}_w{int(pos['width_percent'])}_h{int(pos['height_percent'])}_type{type_name}"


def change_text_to(shape, new_text: str, font_color: tuple = None) -> None:
    """
    Change text in a shape while preserving font style.

    Args:
        shape: The PowerPoint shape to modify
        new_text: The new text to set
        font_color: Optional tuple of (r,g,b) values for font color
    """
    try:
        text_frame = shape.text_frame
        if not text_frame:
            logger.warning("Shape has no text frame")
            return

        if not text_frame.paragraphs:
            logger.warning("Shape has no paragraphs")
            return

        # 원본 스타일 저장
        font_props = {}
        font_found = False

        for paragraph in text_frame.paragraphs:
            if font_found:
                break
            for run in paragraph.runs:
                if run.text and run.text.strip():
                    font_props['name'] = run.font.name
                    font_props['size'] = run.font.size
                    font_props['bold'] = run.font.bold
                    font_props['italic'] = run.font.italic
                    font_props['color'] = safe_get_font_color(run)
                    font_found = True
                    break

        # 폰트 정보 없을 경우도 대비
        font_props.setdefault('name', None)
        font_props.setdefault('size', None)
        font_props.setdefault('bold', None)
        font_props.setdefault('italic', None)
        font_props.setdefault('color', None)

        # 텍스트 프레임 초기화
        text_frame.clear()
        new_paragraph = text_frame.paragraphs[0]
        new_run = new_paragraph.add_run()
        new_run.text = new_text

        # 스타일 복원
        if font_props['name']:
            new_run.font.name = font_props['name']
        if font_props['size']:
            new_run.font.size = font_props['size']
        if font_props['bold'] is not None:
            new_run.font.bold = font_props['bold']
        if font_props['italic'] is not None:
            new_run.font.italic = font_props['italic']

        # 폰트 컬러 복원
        try:
            # 매개변수로 전달된 컬러가 있으면 우선 사용
            color_to_use = font_color if font_color is not None else font_props['color']
            
            if color_to_use:
                r, g, b = color_to_use
                new_run.font.color.rgb = RGBColor(r, g, b)
        except Exception as e:
            logger.warning(f"폰트 컬러 설정 중 오류: {e}")

    except Exception as e:
        logger.error(f"Error updating text in shape: {str(e)}")
        raise

def is_tag_identifier(field_name: str) -> bool:
    """필드 이름이 태그 관련 식별자인지 확인합니다."""
    tag_keywords = ["tag", "label", "cic_label", "ui_element"]
    return any(keyword in str(field_name).lower() for keyword in tag_keywords)

def find_tag_element(slide, normalized_field_name: str, slide_width: int, slide_height: int) -> Optional[Any]:
    """슬라이드에서 태그/라벨 요소를 찾습니다."""
    MAX_TAG_LENGTH = 15
    
    # 정확히 일치하는 요소 먼저 검색
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            continue
            
        shape_text = "\n".join([p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()])
        if not shape_text or len(shape_text) > MAX_TAG_LENGTH:
            continue
            
        normalized_shape_text = ''.join(shape_text.lower().split())
        if normalized_shape_text == normalized_field_name:
            return shape
    
    # 정확히 일치하지 않으면 부분 일치 검색 (태그의 경우 tag_1, tag_2와 같은 형태일 수 있음)
    match_pattern = re.compile(r'(tag|label|cic|ui)[_\-\s]?\d*')
    
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            continue
            
        shape_text = "\n".join([p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()])
        if not shape_text or len(shape_text) > MAX_TAG_LENGTH:
            continue
            
        normalized_shape_text = ''.join(shape_text.lower().split())
        
        # 크기가 작은지 확인 (태그/라벨은 일반적으로 작음)
        if hasattr(shape, 'width') and hasattr(shape, 'height'):
            width_percent = (shape.width / slide_width) * 100
            height_percent = (shape.height / slide_height) * 100
            
            is_small = width_percent <= 20 and height_percent <= 10
            
            # 작은 요소이고 패턴이 일치하면 반환
            if is_small and (match_pattern.match(normalized_shape_text) or 
                             match_pattern.match(normalized_field_name)):
                return shape
                
            # 작은 요소이고 두 텍스트가 비슷하면 반환 (편집 거리 활용)
            if is_small and len(normalized_shape_text) <= 10 and len(normalized_field_name) <= 10:
                # 레벤슈타인 거리 계산 (간단 구현)
                def levenshtein_distance(s1, s2):
                    if len(s1) < len(s2):
                        return levenshtein_distance(s2, s1)
                    if len(s2) == 0:
                        return len(s1)
                    previous_row = range(len(s2) + 1)
                    for i, c1 in enumerate(s1):
                        current_row = [i + 1]
                        for j, c2 in enumerate(s2):
                            insertions = previous_row[j + 1] + 1
                            deletions = current_row[j] + 1
                            substitutions = previous_row[j] + (c1 != c2)
                            current_row.append(min(insertions, deletions, substitutions))
                        previous_row = current_row
                    return previous_row[-1]
                
                # 짧은 텍스트에 대해 편집 거리가 작으면 유사하다고 판단
                distance = levenshtein_distance(normalized_shape_text, normalized_field_name)
                if distance <= 3:  # 편집 거리 임계값
                    return shape
    
    return None

def find_shape_by_text_with_count(slide, field_name: str, normalized_field_name: str, expected_count: int = 1) -> Tuple[Optional[Any], int]:
    """텍스트와 카운트를 기준으로 도형을 찾습니다."""
    current_count = 0
    
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            continue
            
        shape_text = "\n".join([p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()])
        if not shape_text:
            continue
            
        normalized_shape_text = ''.join(shape_text.lower().split())
        
        if normalized_shape_text == normalized_field_name or shape_text.lower() == field_name.lower():
            current_count += 1
            
            if current_count == expected_count:
                return shape, current_count
                
    return None, current_count

def extract_count_from_field_name(field_name: str) -> Tuple[str, str, int]:
    """필드 이름에서 카운트 정보를 추출합니다."""
    count_match = re.search(r'_(\d+)$', field_name)
    expected_count = 1
    
    if count_match:
        try:
            expected_count = int(count_match.group(1))
            clean_field_name = field_name[:count_match.start()]
            clean_normalized_field_name = ''.join(clean_field_name.lower().split())
            return clean_field_name, clean_normalized_field_name, expected_count
        except (ValueError, IndexError):
            pass
    
    return field_name, ''.join(field_name.lower().split()), expected_count

def is_special_content(text: str) -> bool:
    """텍스트가 숫자나 특수 기호로만 구성되어 있는지 확인합니다."""
    # 공백 제거
    text = text.strip()
    if not text:
        return False
    
    # 숫자로만 구성된 경우
    if text.isdigit():
        return True

    if len(text) < 3: 
        return True

    # 특수 기호 목록 (추가된 기호들)
    special_chars = set('!@#$%^&*()_+-=[]{}|\\;:\'",.<>/?`~・×÷©®™℠℗℮℅℆℄℀℁℃℉№℡™')
    
    # 텍스트가 특수 기호로만 구성된 경우
    if all(c in special_chars or c.isspace() for c in text):
        return True
    
    # 숫자와 특수 기호의 조합인 경우
    if not any(c.isalpha() or '\u3131' <= c <= '\u318E' or '\uAC00' <= c <= '\uD7A3' for c in text):
        return True
    
    return False

def generate_position_key(text: str, position: dict) -> str:
    """텍스트와 위치 정보를 조합하여 고유 키 생성"""
    # 특수 문자나 숫자만 있는 경우 텍스트 그대로 반환
    if is_special_content(text):
        return text
        
    # 위치를 5% 단위로 반올림하여 근접 위치는 같은 키로 처리
    left_group = round(position['left_percent'] / 5) * 5
    top_group = round(position['top_percent'] / 5) * 5
    return f"{text}__pos_{left_group}_{top_group}"


def load_presentation(template_path: str) -> Presentation:
    """프레젠테이션 파일을 로드합니다."""
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    try:
        return Presentation(template_path)
    except Exception as e:
        logger.error(f"프레젠테이션 로드 중 오류 발생: {str(e)}")
        raise

def save_presentation(prs: Presentation, output_path: str) -> None:
    """프레젠테이션을 안전하게 저장합니다."""
    try:
        # # 기존 파일이 있다면 삭제
        # if os.path.exists(output_path):
        #     os.remove(output_path)
        #     logger.info(f"기존 파일 삭제: {output_path}")
        
        # 새로운 파일 저장
        prs.save(output_path)
        logger.info(f"프레젠테이션 저장 완료: {output_path}")
    except Exception as e:
        logger.error(f"프레젠테이션 저장 중 오류 발생: {str(e)}")
        raise

def save_meta_info(meta_data: dict, meta_path: str) -> None:
    """메타 정보를 안전하게 저장합니다."""
    try:
        # 기존 파일이 있다면 삭제
        if os.path.exists(meta_path):
            os.remove(meta_path)
            logger.info(f"기존 메타 파일 삭제: {meta_path}")
        
        # 새로운 파일 저장
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(meta_data, f, ensure_ascii=False, indent=2)
        logger.info(f"메타 정보 저장 완료: {meta_path}")
    except Exception as e:
        logger.error(f"메타 정보 저장 중 오류 발생: {str(e)}")
        raise

def load_meta_info(meta_path: str) -> Dict[str, Any]:
    """메타 정보를 로드합니다."""
    if not os.path.exists(meta_path):
        raise FileNotFoundError(f"Meta file not found: {meta_path}")
    try:
        with open(meta_path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception as e:
        logger.error(f"메타 정보 로드 중 오류 발생: {str(e)}")
        raise

def ensure_directory(directory_path: str) -> None:
    """디렉토리가 존재하지 않으면 생성합니다."""
    try:
        os.makedirs(directory_path, exist_ok=True)
        logger.info(f"디렉토리 확인/생성 완료: {directory_path}")
    except Exception as e:
        logger.error(f"디렉토리 생성 중 오류 발생: {str(e)}")
        raise

def extract_text_and_style(shape) -> Dict[str, Any]:
    """도형에서 텍스트와 스타일 정보를 추출합니다."""
    result = {
        'text': '',
        'style': None,
        'has_text': False
    }

    if not hasattr(shape, 'text_frame') or not shape.text_frame:
        return result

    text_frame = shape.text_frame
    text_parts = []
    first_text_style = None

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text and run.text.strip():
                text_parts.append(run.text)

                if first_text_style is None:
                    first_text_style = {
                        'font_name': run.font.name,
                        'font_size': run.font.size,
                        'font_bold': run.font.bold,
                        'font_italic': run.font.italic,
                        'font_color': safe_get_font_color(run)
                    }

    combined_text = ' '.join(text_parts).strip()
    result['text'] = combined_text
    result['style'] = first_text_style
    result['has_text'] = bool(combined_text)

    return result

def create_unique_text_key(text: str, position: Dict[str, float], slide_number: int) -> str:
    """텍스트와 위치 정보를 조합하여 unique key를 생성합니다."""
    position_str = f"s{slide_number}_l{int(position['left_percent'])}_t{int(position['top_percent'])}"
    return f"{text}_{position_str}"


def process_shape_for_meta(shape, slide_number: int, slide_width: int, slide_height: int) -> Dict[str, Any]:
    """도형의 메타 정보를 생성합니다."""
    # 그룹 도형 처리
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group_shapes = extract_group_shapes(shape)
        group_elements = []
        
        for shape_info in group_shapes:
            position = get_shape_position(shape_info['shape'], slide_width, slide_height)
            group_elements.append({
                'text': shape_info['text_info']['text'],
                'position': position,
                'styles': shape_info['text_info']['styles'],
                'paragraph_styles': shape_info['text_info']['paragraph_styles'],
                'type': shape_info['type'],
                'unique_key': create_unique_text_key(shape_info['text_info']['text'], position, slide_number)
            })
        
        if group_elements:
            return {
                'type': 'GROUP',
                'elements': group_elements,
                'position': get_shape_position(shape, slide_width, slide_height)
            }
        return None
    
    # 일반 도형 처리 (기존 로직)
    position = get_shape_position(shape, slide_width, slide_height)
    text_info = extract_text_and_style(shape)
    
    if not text_info['has_text']:
        return None
        
    return {
        'text': text_info['text'],
        'position': position,
        'styles': text_info['styles'],
        'paragraph_styles': text_info['paragraph_styles'],
        'type': get_type_info(shape),
        'unique_key': create_unique_text_key(text_info['text'], position, slide_number)
    }

def extract_presentation_meta(prs: Presentation) -> Dict[str, Any]:
    """프레젠테이션의 모든 메타 정보를 추출합니다."""
    meta_data = {
        'slides': [],
        'styles': {},
        'text_elements': {}
    }
    
    for slide_number, slide in enumerate(prs.slides, 1):
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        slide_elements = []
        for shape in slide.shapes:
            shape_meta = process_shape_for_meta(shape, slide_number, slide_width, slide_height)
            if shape_meta:
                slide_elements.append(shape_meta)
                # 스타일 정보 저장
                if shape_meta['style']:
                    meta_data['styles'][shape_meta['unique_key']] = shape_meta['style']
                # 텍스트 요소 정보 저장
                meta_data['text_elements'][shape_meta['unique_key']] = {
                    'text': shape_meta['text'],
                    'position': shape_meta['position'],
                    'type': shape_meta['type']
                }
        
        meta_data['slides'].append({
            'slide_number': slide_number,
            'elements': slide_elements
        })
    
    return meta_data

def apply_text_roles(prs: Presentation, role_mapping: Dict[str, str]) -> None:
    """텍스트 요소에 role을 적용합니다."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue
                
            position = get_shape_position(shape, prs.slide_width, prs.slide_height)
            text_info = extract_text_and_style(shape)
            
            if not text_info['has_text']:
                continue
                
            unique_key = create_unique_text_key(
                text_info['text'],
                position,
                slide.slide_id  # slide number로 대체 가능
            )
            
            if unique_key in role_mapping:
                change_text_to(shape, role_mapping[unique_key])

def extract_group_shapes(shape) -> List[Dict[str, Any]]:
    """그룹 도형에서 텍스트를 포함한 모든 하위 도형을 추출합니다."""
    shapes_info = []
    
    def process_shape_in_group(shape):
        # 그룹 도형인 경우 재귀적으로 처리
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                process_shape_in_group(child)
        else:
            # 텍스트가 있는 도형만 처리
            text_info = extract_text_and_style(shape)
            if text_info['has_text']:
                shape_info = {
                    'type': get_type_info(shape),
                    'text_info': text_info,
                    'shape': shape  # 원본 도형 객체 참조 저장
                }
                shapes_info.append(shape_info)
    
    # 최상위 도형이 그룹인 경우에만 처리
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            process_shape_in_group(child)
    
    return shapes_info



def save_single_slide(slide, output_path: str) -> None:
    """
    단일 슬라이드를 새로운 프레젠테이션 파일로 저장합니다.
    
    Args:
        slide: 저장할 슬라이드 객체
        output_path: 저장할 파일 경로
    """
    try:
        # 새 프레젠테이션 객체 생성
        prs = Presentation()
        
        # 원본 슬라이드 복사
        slide_layout = prs.slide_layouts[0]  # 기본 레이아웃 사용
        new_slide = prs.slides.add_slide(slide_layout)
        
        # 원본 슬라이드의 모든 도형 복사
        for shape in slide.shapes:
            el = shape.element
            new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
        
        # 결과 저장
        prs.save(str(output_path))
        logger.info(f"슬라이드를 성공적으로 저장했습니다: {output_path}")
        
    except Exception as e:
        logger.error(f"슬라이드 저장 중 오류 발생: {str(e)}")
        raise 
