import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging
from typing import Dict, Any, Optional

# 공통 모듈에서 함수 가져오기
from ppt_common import (
    get_shape_position,
    change_text_to, 
    is_tag_identifier, 
    find_tag_element,
    find_shape_by_text_with_count,
    extract_count_from_field_name,
    generate_position_key,
    safe_get_font_color, 
    logger,
    # 파일 처리 관련 함수 추가
    load_presentation,
    save_presentation,
    save_meta_info,
    load_meta_info,
    ensure_directory
)

def find_shape_by_position_key(slide, position_key: str, slide_width: int, slide_height: int) -> Optional[Any]:
    """위치 기반 키를 사용하여 도형 찾기"""
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
            
        # 각 문단의 텍스트를 추출하고 줄바꿈으로 합치기
        paragraphs_text = []
        for paragraph in shape.text_frame.paragraphs:
            paragraph_runs = []
            for run in paragraph.runs:
                if run.text and run.text.strip():
                    paragraph_runs.append(run.text.strip())
            if paragraph_runs:
                paragraphs_text.append(" ".join(paragraph_runs))
        
        shape_text = "\n".join(paragraphs_text)
        if not shape_text:
            continue
            
        # 현재 도형의 위치 정보로 키 생성
        pos = get_shape_position(shape, slide_width, slide_height)
        current_key = generate_position_key(shape_text, pos)
        
        # 키가 일치하면 해당 도형 반환
        if current_key == position_key.split('_', 1)[0]:  # _1, _2 등의 suffix 제거
            return shape
            
    return None

def process_tag_element(slide, field_name, content, normalized_field_name, slide_width, slide_height, font_color):
    """
    태그/라벨 요소를 처리합니다.
    
    Args:
        slide: 처리할 슬라이드
        field_name: 필드 이름
        content: 업데이트할 내용
        normalized_field_name: 정규화된 필드 이름
        slide_width: 슬라이드 너비
        slide_height: 슬라이드 높이
        font_color: 폰트 컬러
        
    Returns:
        bool: 처리 성공 여부
    """
    # 태그/라벨 요소 찾기
    tag_shape = find_tag_element(slide, normalized_field_name, slide_width, slide_height)
    
    if tag_shape:
        logger.info(f"태그/라벨 요소 발견: '{field_name}' -> '{content}'")
        meta_info = {
            'styles': [],
            'paragraph_styles': [],
            'font_color': font_color
        }
        change_text_to(tag_shape, str(content), meta_info)
        return True
        
    return False

def update_table_cell(slide, field_name: str, content: str, table_info: dict, normalized_field_name: str, font_color) -> bool:
    """표 셀의 텍스트를 업데이트합니다."""
    row = table_info.get('row', 0)
    col = table_info.get('col', 0)
    
    # 필드 이름에서 카운트 추출
    clean_field_name, clean_normalized_field_name, expected_count = extract_count_from_field_name(field_name)
    
    # 메타 정보에서 text_count 가져오기
    meta_info = getattr(update_table_cell, 'meta_info', {})
    if meta_info and field_name in meta_info:
        text_count = meta_info[field_name].get('text_count')
        if text_count:
            expected_count = text_count
    
    # 표 인스턴스 카운터
    current_count = 0
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                # 표 행렬이 충분히 큰지 확인
                if len(shape.table.rows) >= row and len(shape.table.columns) >= col:
                    # 1부터 시작하므로 인덱스는 1 빼기
                    cell = shape.table.cell(row-1, col-1)
                    cell_text = "\n".join([p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip()])
                    
                    if not cell_text:
                        continue
                        
                    # 셀 텍스트 정규화
                    normalized_cell_text = ''.join(cell_text.lower().split())
                    
                    # 텍스트 일치 확인
                    if normalized_cell_text == clean_normalized_field_name or cell_text.lower() == clean_field_name.lower():
                        current_count += 1
                        
                        # 원하는 순번의 인스턴스를 찾았는지 확인
                        if current_count == expected_count:
                            logger.info(f"표 셀 일치 발견 [행:{row}, 열:{col}, 인스턴스 #{expected_count}]: '{cell_text[:30]}...' -> '{content[:30]}...'")
                            meta_info = {
                                'styles': [],
                                'paragraph_styles': [],
                                'font_color': font_color
                            }
                            change_text_to(cell, str(content), meta_info)
                            return True
            except Exception as e:
                logger.error(f"표 셀 업데이트 중 오류: {e}")
    
    if current_count > 0:
        logger.info(f"경고: 요청한 표 셀을 {current_count}개 찾았지만, 요청한 {expected_count}번째 인스턴스는 찾지 못했습니다.")
    
    return False

def process_regular_shapes(slide, field_name: str, content: str, normalized_field_name: str, font_color, expected_count=1) -> bool:
    """일반 도형들의 텍스트를 처리합니다."""
    # 필드 이름에서 카운트 추출
    clean_field_name, clean_normalized_field_name, count_from_name = extract_count_from_field_name(field_name)
    
    # 명시적으로 지정된 expected_count가 없으면 필드 이름에서 추출한 count 사용
    if expected_count == 1 and count_from_name > 1:
        expected_count = count_from_name
    
    # 메타 정보에서 text_count 가져오기
    meta_info = getattr(process_regular_shapes, 'meta_info', {})
    if meta_info and field_name in meta_info:
        text_count = meta_info[field_name].get('text_count')
        if text_count:
            expected_count = text_count
    
    # 번호 기반으로 해당 요소 찾기
    shape, found_count = find_shape_by_text_with_count(slide, clean_field_name, clean_normalized_field_name, expected_count)
    
    if shape:
        logger.info(f"텍스트 일치 발견 (인스턴스 #{expected_count}): '{clean_field_name[:30]}...' -> '{content[:30]}...'")
        meta_info = {
            'styles': [],
            'paragraph_styles': [],
            'font_color': font_color
        }
        change_text_to(shape, str(content), meta_info)
        return True
            
    # 여러 인스턴스 중 하나를 못 찾았어도 다른 인스턴스가 있는지 확인
    if found_count > 0:
        logger.info(f"경고: '{clean_field_name}' 텍스트를 가진 요소를 {found_count}개 찾았지만, 요청한 {expected_count}번째 인스턴스는 찾지 못했습니다.")
        
    return False

def process_group_shapes(slide, field_name: str, content: str, normalized_field_name: str, font_color) -> bool:
    """그룹 요소 내부의 도형들을 처리합니다. 중첩 그룹도 처리합니다."""
    found_match = False
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                # 그룹 내부 처리를 재귀적으로 수행하는 내부 함수
                found_match = process_group_recursive(shape, field_name, content, normalized_field_name, font_color)
                if found_match:
                    return True
            except Exception as e:
                logger.error(f"그룹 요소 처리 중 오류: {e}")
    
    return False

def process_group_recursive(group_shape, field_name: str, content: str, normalized_field_name: str, font_color, depth: int = 1, expected_count: int = 1) -> bool:
    """그룹 내부를 재귀적으로 처리하는 함수"""
    # 필드 이름에서 카운트 추출
    clean_field_name, clean_normalized_field_name, count_from_name = extract_count_from_field_name(field_name)
    
    # 명시적으로 지정된 expected_count가 없으면 필드 이름에서 추출한 count 사용
    if expected_count == 1 and count_from_name > 1:
        expected_count = count_from_name
    
    # 현재 그룹 내 일치하는 요소 카운터
    current_count = 0
    
    # 그룹 내의 모든 도형을 처리
    for shape_in_group in group_shape.shapes:
        # 중첩 그룹 처리
        if shape_in_group.shape_type == MSO_SHAPE_TYPE.GROUP:
            logger.info(f"중첩 그룹 발견 (깊이: {depth})")
            found = process_group_recursive(shape_in_group, clean_field_name, content, clean_normalized_field_name, font_color, depth + 1, expected_count)
            if found:
                return True
        
        # 텍스트 프레임이 있는 도형 처리
        if hasattr(shape_in_group, "text_frame") and shape_in_group.text_frame:
            # 그룹 내 도형 텍스트 추출
            shape_text = "\n".join([p.text.strip() for p in shape_in_group.text_frame.paragraphs if p.text.strip()])
            
            if not shape_text:
                continue
            
            # 텍스트 정규화
            normalized_shape_text = ''.join(shape_text.lower().split())
            
            # 텍스트 비교 - 정규화된 텍스트 또는 원본 텍스트 일치 확인
            if normalized_shape_text == clean_normalized_field_name or shape_text.lower() == clean_field_name.lower():
                current_count += 1
                
                # 원하는 순번의 인스턴스 발견
                if current_count == expected_count:
                    logger.info(f"그룹 내 텍스트 일치 발견 (깊이: {depth}, 인스턴스 #{expected_count}): '{shape_text[:30]}...' -> '{content[:30]}...'")
                    meta_info = {
                        'styles': [],
                        'paragraph_styles': [],
                        'font_color': font_color
                    }
                    change_text_to(shape_in_group, str(content), meta_info)
                    return True
    
    # 여러 인스턴스 중 하나를 못 찾았어도 다른 인스턴스가 있는지 확인
    if current_count > 0:
        logger.info(f"경고: 그룹 내에서 '{clean_field_name}' 텍스트를 가진 요소를 {current_count}개 찾았지만, 요청한 {expected_count}번째 인스턴스는 찾지 못했습니다.")
    
    return False

def extract_text_from_shape(shape) -> dict:
    """도형에서 텍스트와 폰트 컬러를 추출합니다."""
    if not hasattr(shape, "text_frame"):
        return {"text": "", "font_color": None}
    
    text_parts = []
    font_colors = []
    
    for paragraph in shape.text_frame.paragraphs:
        paragraph_runs = []
        paragraph_colors = []
        
        for run in paragraph.runs:
            if run.text and run.text.strip():
                paragraph_runs.append(run.text.strip())
                # 폰트 컬러 추출 개선
                try:
                    # 1. 직접 font.color.rgb 확인
                    if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                        paragraph_colors.append((
                            run.font.color.rgb.r,
                            run.font.color.rgb.g,
                            run.font.color.rgb.b
                        ))
                    # 2. 도형의 fill 속성 확인
                    elif hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                        paragraph_colors.append((
                            shape.fill.fore_color.rgb.r,
                            shape.fill.fore_color.rgb.g,
                            shape.fill.fore_color.rgb.b
                        ))
                    else:
                        paragraph_colors.append(None)
                except Exception as e:
                    logger.debug(f"폰트 컬러 추출 중 오류 발생: {e}")
                    paragraph_colors.append(None)
        
        if paragraph_runs:
            text_parts.append(" ".join(paragraph_runs))
            # 단락의 대표 컬러 선택 (None이 아닌 첫 번째 컬러)
            valid_colors = [c for c in paragraph_colors if c is not None]
            font_colors.append(valid_colors[0] if valid_colors else None)
    
    # 전체 텍스트의 대표 컬러 선택
    valid_colors = [c for c in font_colors if c is not None]
    final_color = valid_colors[0] if valid_colors else (0, 0, 0)  # 기본값은 검정색
    
    return {
        "text": "\n".join(text_parts),
        "font_color": final_color
    }

def update_slide(slide, schema, slide_width, slide_height):
    """슬라이드 내용을 업데이트합니다."""
    logger.info(f"슬라이드 업데이트 시작: {schema['fields']}")
    
    for field_key, field_info in schema['fields'].items():
        # role 값이 없으면 건너뛰기
        if 'role' not in field_info:
            logger.warning(f"필드 '{field_key}'에 role 정보가 없습니다.")
            continue
            
        content = field_info['role']  # role 값을 새로운 텍스트로 사용
        
        # 스타일 정보 구성
        meta_info = {
            'font_styles': field_info.get('font_styles', []),  # 폰트 스타일 정보
            'font_color': field_info.get('font_color'),  # 기본 폰트 컬러
            'paragraph_styles': field_info.get('paragraph_styles', [])  # 기존 단락 스타일 유지
        }
        
        # 위치 기반 키로 도형 찾기
        shape = find_shape_by_position_key(slide, field_key, slide_width, slide_height)
        if shape:
            logger.info(f"위치 기반 매칭 성공: '{field_key}' -> '{content}'")
            # 스타일 정보와 함께 텍스트 업데이트
            if not change_text_to(shape, str(content), meta_info):
                logger.warning(f"텍스트 업데이트 실패: {field_key}")
            continue
            
        # 위치 기반 매칭 실패 시 기존 방식으로 시도
        original_text = field_info.get('original_text', '')
        normalized_field_name = ''.join(original_text.lower().split())
        
        # 1. 태그/라벨 요소 처리
        if is_tag_identifier(original_text) or is_tag_identifier(content):
            if process_tag_element(slide, original_text, content, normalized_field_name, slide_width, slide_height, meta_info):
                continue
        
        # 2. 표 셀 처리
        table_info = field_info.get('table_info', None)
        if table_info:
            if update_table_cell(slide, original_text, content, table_info, normalized_field_name, meta_info):
                continue
        
        # 3. 일반 도형 처리
        text_count = field_info.get('text_count', 1)
        if process_regular_shapes(slide, original_text, content, normalized_field_name, meta_info, text_count):
            continue
            
        # 4. 그룹 요소 처리
        if not process_group_shapes(slide, original_text, content, normalized_field_name, meta_info):
            logger.warning(f"필드 '{field_key}'에 대한 매칭 요소를 찾지 못했습니다.")
