import os
import json
from pptx import Presentation
from langchain_openai import AzureChatOpenAI
from langchain.prompts import ChatPromptTemplate
from dotenv import load_dotenv 

load_dotenv()

def extract_text_items(pptx_path):
    prs = Presentation(pptx_path)
    text_items = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip():
                text_items.append({
                    "slide": slide_idx,
                    "type": "textbox",
                    "text": shape.text_frame.text.strip()
                })
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows, 1):
                    for col_idx, cell in enumerate(row.cells, 1):
                        if cell.text_frame and cell.text_frame.text.strip():
                            text_items.append({
                                "slide": slide_idx,
                                "type": "table_cell",
                                "row": row_idx,
                                "col": col_idx,
                                "text": cell.text_frame.text.strip()
                            })
    return text_items

# 랭체인 LLM 객체 생성
llm = AzureChatOpenAI(
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_key=os.getenv("AZURE_OPENAI_API_KEY"),
    azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT"), 
    api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2023-05-15"),
    # temperature=0.3,
    # max_tokens=50,
)



def get_ppt_goal(text_items):
    all_text = "\n".join([f"슬라이드 {item['slide']}: {item['text']}" for item in text_items])
    prompt = ChatPromptTemplate.from_template(
        "아래는 파워포인트의 모든 텍스트입니다.\n---\n{all_text}\n---\n이 PPT가 전체적으로 추구하는 목표나 의도를 한 줄로 요약해줘."
    )
    messages = prompt.format_messages(all_text=all_text)
    return llm.invoke(messages).content.strip()

def generate_template_mapping(text_items, goal_summary):
    mapping = {}
    prompt_template = ChatPromptTemplate.from_template(
        "아래는 파워포인트의 일부 텍스트입니다.\n"
        "- 전체 목표: \"{goal_summary}\"\n"
        "- 위치: 슬라이드 {slide}, {type}, {row}, {col}\n"
        "- 원본 텍스트: \"{text}\"\n\n"
        "이 텍스트가 템플릿에서 어떤 역할을 하는지, 아래 예시처럼 영문 snake_case로 식별자만 반환해줘.\n"
        "반드시 아래 예시처럼 텍스트의 의미와 위치에 따라 구조적으로 매핑해줘.\n\n"
        "예시)\n"
        "- \"0. 들어가며\" → title\n"
        "- \"구조개선과 더불어 ...\" → one_line_summary\n"
        "- \"UP가전, 스마트가전 실현을 위한 필요 활동\" → subtitle\n"
        "- \"구조개선\" → category1\n"
        "- \"UP 단위 별 개발과 배포가...\" → category1_subtitle\n"
        "- \"제품 SW플랫폼 구축\" → category1_activity_detail\n"
        "- \"개선 목표에 따른 SW 재구조화...\" → category1_activity_highlight\n"
        "- \"상시 검증\" → category2\n"
        "- \"UP 단위 별 검증과...\" → category2_subtitle\n"
        "- \"레퍼런스 테스트 환경...\" → category2_activity_highlight\n\n"
        "위 예시를 참고해서, 반드시 영문 snake_case로 식별자만 반환해줘."
    )
    for item in text_items:
        max_len = len(item['text'])
        messages = prompt_template.format_messages(
            goal_summary=goal_summary,
            slide=item['slide'],
            type=item.get('type', ''),
            row=item.get('row', ''),
            col=item.get('col', ''),
            text=item['text'],
            max_len=max_len
        )
        template_text = llm.invoke(messages).content.strip()
        mapping[item['text']] = {"text": template_text}
    return mapping

if __name__ == "__main__":
    pptx_path = "/home/jayseo/workspace/hakkamakka/data/input/template_simple.pptx"
    text_items = extract_text_items(pptx_path)
    goal_summary = get_ppt_goal(text_items)
    mapping = generate_template_mapping(text_items, goal_summary)
    with open("/home/jayseo/workspace/hakkamakka/data/output/template_simple_mapping.json", "w", encoding="utf-8") as f:
        json.dump(mapping, f, ensure_ascii=False, indent=2)
    print("매핑 테이블이 /home/jayseo/workspace/hakkamakka/data/output/template_simple_mapping.json에 저장되었습니다.") 
